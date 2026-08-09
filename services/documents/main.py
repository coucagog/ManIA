"""
Service de documents partagé MANIA — conversion + remplissage de gabarits.

Moteurs :
  - Conversion : LibreOffice (soffice), file sérialisée, profil isolé (§7).
  - Remplissage : docxtpl (docx, Jinja2 complet, EN BAC À SABLE), openpyxl (xlsx,
    substitution {{clé}}), python-pptx (pptx, substitution {{clé}}).

Règles (STACK.md §7, STACK-3.md §36/§42/§44) :
  - Conteneur SÉPARÉ de la transcription.
  - Conversion soffice sérialisée (sémaphore=1) — jamais 2 LibreOffice à la fois.
  - Éphémère : dossier temporaire par requête, supprimé ; rien stocké ni loggué.
  - Token par tenant dérivé (partagé avec la transcription).

Endpoints :
  POST /v1/convert   fichier + `to`                   -> fichier converti.
  POST /v1/fill      gabarit + `data` (JSON) [+ `to`] -> document rempli (+ converti).
  GET  /health       (sans auth)
"""

from __future__ import annotations

import asyncio
import hashlib
import hmac
import io
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
import time
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, File, Form, Header, HTTPException, UploadFile
from fastapi.responses import StreamingResponse

SHARED_SERVICES_SECRET = os.environ.get("SHARED_SERVICES_SECRET", "")
MAX_UPLOAD_MB = int(os.environ.get("MAX_UPLOAD_MB", "50"))
CONVERT_TIMEOUT_S = int(os.environ.get("CONVERT_TIMEOUT_S", "120"))
SOFFICE_BIN = os.environ.get("SOFFICE_BIN", "soffice")

ALLOWED_FORMATS = {
    "pdf", "docx", "odt", "doc", "rtf", "txt", "html",
    "xlsx", "ods", "csv", "pptx", "odp", "ppt", "xls",
    "png", "jpg", "svg", "epub",
}
FILL_FORMATS = {"docx", "xlsx", "pptx"}

MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc": "application/msword",
    "odt": "application/vnd.oasis.opendocument.text",
    "rtf": "application/rtf",
    "txt": "text/plain",
    "html": "text/html",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "xls": "application/vnd.ms-excel",
    "ods": "application/vnd.oasis.opendocument.spreadsheet",
    "csv": "text/csv",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt": "application/vnd.ms-powerpoint",
    "odp": "application/vnd.oasis.opendocument.presentation",
    "png": "image/png",
    "jpg": "image/jpeg",
    "svg": "image/svg+xml",
    "epub": "application/epub+zip",
}

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger("mania-documents")

_semaphore: Optional[asyncio.Semaphore] = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    global _semaphore
    if not SHARED_SERVICES_SECRET:
        raise RuntimeError("SHARED_SERVICES_SECRET manquant — refus de démarrer sans lui.")
    _semaphore = asyncio.Semaphore(1)
    try:
        subprocess.run([SOFFICE_BIN, "--headless", "--version"], capture_output=True, timeout=60)
        log.info("soffice prêt.")
    except Exception:
        log.warning("pré-chauffage soffice impossible (on continue).")
    yield


app = FastAPI(title="mania-documents", lifespan=lifespan, docs_url=None, redoc_url=None, openapi_url=None)


def verify_token(authorization: Optional[str]) -> str:
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Authorization Bearer manquant")
    token = authorization[len("Bearer "):].strip()
    if "." not in token:
        raise HTTPException(status_code=401, detail="Token malformé")
    slug, _, presented_hex = token.rpartition(".")
    if not slug or not presented_hex:
        raise HTTPException(status_code=401, detail="Token malformé")
    expected_hex = hmac.new(SHARED_SERVICES_SECRET.encode("utf-8"), slug.encode("utf-8"), hashlib.sha256).hexdigest()
    if not hmac.compare_digest(expected_hex, presented_hex):
        raise HTTPException(status_code=401, detail="Token invalide")
    return slug


@app.get("/health")
async def health():
    return {
        "status": "ok",
        "engine": "libreoffice",
        "convert_formats": sorted(ALLOWED_FORMATS),
        "fill_templates": sorted(FILL_FORMATS),
    }


def _run_soffice(workdir: str, in_path: str, fmt: str) -> str:
    profile = Path(workdir) / "profile"
    outdir = Path(workdir) / "out"
    outdir.mkdir(parents=True, exist_ok=True)
    proc = subprocess.run(
        [
            SOFFICE_BIN,
            f"-env:UserInstallation={profile.as_uri()}",
            "--headless", "--norestore", "--nolockcheck",
            "--convert-to", fmt,
            "--outdir", str(outdir),
            in_path,
        ],
        capture_output=True, timeout=CONVERT_TIMEOUT_S,
    )
    produced = list(outdir.glob(f"*.{fmt}"))
    if not produced:
        err = (proc.stderr or b"").decode("utf-8", "replace")[:200]
        raise RuntimeError(f"aucune sortie .{fmt} (soffice rc={proc.returncode}) {err}")
    return str(produced[0])


async def _convert_serialized(workdir: str, in_path: str, fmt: str) -> str:
    assert _semaphore is not None
    loop = asyncio.get_running_loop()
    async with _semaphore:
        return await asyncio.wait_for(
            loop.run_in_executor(None, _run_soffice, workdir, in_path, fmt),
            timeout=CONVERT_TIMEOUT_S + 15,
        )


_PLACEHOLDER = re.compile(r"\{\{\s*([\w.]+)\s*\}\}")


def _substitute(text: str, ctx: dict) -> str:
    def repl(m):
        key = m.group(1).strip()
        return str(ctx[key]) if key in ctx else m.group(0)
    return _PLACEHOLDER.sub(repl, text)


def _fill_docx(tpl: str, out: str, ctx: dict) -> None:
    from docxtpl import DocxTemplate
    from jinja2.sandbox import SandboxedEnvironment
    doc = DocxTemplate(tpl)
    doc.render(ctx, jinja_env=SandboxedEnvironment())
    doc.save(out)


def _fill_xlsx(tpl: str, out: str, ctx: dict) -> None:
    from openpyxl import load_workbook
    wb = load_workbook(tpl)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and "{{" in cell.value:
                    cell.value = _substitute(cell.value, ctx)
    wb.save(out)


def _fill_pptx(tpl: str, out: str, ctx: dict) -> None:
    from pptx import Presentation
    prs = Presentation(tpl)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and "{{" in run.text:
                        run.text = _substitute(run.text, ctx)
    prs.save(out)


def _fill_template(ext: str, tpl: str, out: str, ctx: dict) -> None:
    if ext == "docx":
        _fill_docx(tpl, out, ctx)
    elif ext == "xlsx":
        _fill_xlsx(tpl, out, ctx)
    elif ext == "pptx":
        _fill_pptx(tpl, out, ctx)
    else:
        raise ValueError(f"type de gabarit non géré: {ext}")


@app.post("/v1/convert")
async def convert(
    file: UploadFile = File(...),
    to: str = Form(...),
    authorization: Optional[str] = Header(default=None),
):
    tenant = verify_token(authorization)
    fmt = to.lower().strip().lstrip(".")
    if fmt not in ALLOWED_FORMATS:
        raise HTTPException(status_code=400, detail=f"format de sortie non supporté: '{fmt}'")

    raw = await file.read()
    size_mb = len(raw) / (1024 * 1024)
    if size_mb > MAX_UPLOAD_MB:
        raise HTTPException(status_code=413, detail=f"fichier trop volumineux ({size_mb:.1f} Mo > {MAX_UPLOAD_MB} Mo)")

    in_name = Path(file.filename or "input").name or "input"
    stem = Path(in_name).stem or "document"
    workdir = tempfile.mkdtemp(prefix="mania-doc-")
    started = time.monotonic()
    try:
        in_path = os.path.join(workdir, in_name)
        with open(in_path, "wb") as f:
            f.write(raw)
        del raw
        try:
            out_path = await _convert_serialized(workdir, in_path, fmt)
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="conversion trop longue")
        except subprocess.TimeoutExpired:
            raise HTTPException(status_code=504, detail="conversion trop longue")
        except RuntimeError:
            log.warning("tenant=%s conversion échouée -> %s", tenant, fmt)
            raise HTTPException(status_code=422, detail="conversion échouée (format d'entrée non géré, ou fichier invalide)")
        except HTTPException:
            raise
        except Exception:
            log.exception("tenant=%s erreur conversion", tenant)
            raise HTTPException(status_code=500, detail="erreur de conversion")
        data_out = open(out_path, "rb").read()
    finally:
        shutil.rmtree(workdir, ignore_errors=True)

    log.info("tenant=%s convert -> %s en %.1fs (%.1f Mo)", tenant, fmt, time.monotonic() - started, size_mb)
    return StreamingResponse(
        io.BytesIO(data_out),
        media_type=MIME.get(fmt, "application/octet-stream"),
        headers={"Content-Disposition": f'attachment; filename="{stem}.{fmt}"'},
    )


@app.post("/v1/fill")
async def fill(
    template: UploadFile = File(...),
    data: str = Form(...),
    to: Optional[str] = Form(default=None),
    authorization: Optional[str] = Header(default=None),
):
    tenant = verify_token(authorization)

    try:
        ctx = json.loads(data)
    except Exception:
        raise HTTPException(status_code=400, detail="champ 'data' invalide (JSON attendu)")
    if not isinstance(ctx, dict):
        raise HTTPException(status_code=400, detail="champ 'data' doit être un objet JSON")

    tname = Path(template.filename or "modele").name or "modele"
    ext = Path(tname).suffix.lower().lstrip(".")
    if ext not in FILL_FORMATS:
        raise HTTPException(status_code=400, detail=f"gabarit non supporté: '.{ext}' (docx, xlsx ou pptx)")

    target = (to.lower().strip().lstrip(".") if to else ext)
    if target not in ALLOWED_FORMATS:
        raise HTTPException(status_code=400, detail=f"format de sortie non supporté: '{target}'")

    raw = await template.read()
    size_mb = len(raw) / (1024 * 1024)
    if size_mb > MAX_UPLOAD_MB:
        raise HTTPException(status_code=413, detail=f"gabarit trop volumineux ({size_mb:.1f} Mo > {MAX_UPLOAD_MB} Mo)")

    stem = Path(tname).stem or "document"
    workdir = tempfile.mkdtemp(prefix="mania-fill-")
    started = time.monotonic()
    try:
        tpl_path = os.path.join(workdir, tname)
        with open(tpl_path, "wb") as f:
            f.write(raw)
        del raw

        filled_path = os.path.join(workdir, f"rempli.{ext}")
        loop = asyncio.get_running_loop()

        try:
            await asyncio.wait_for(
                loop.run_in_executor(None, _fill_template, ext, tpl_path, filled_path, ctx),
                timeout=CONVERT_TIMEOUT_S,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="remplissage trop long")
        except HTTPException:
            raise
        except Exception:
            log.warning("tenant=%s remplissage échoué (%s)", tenant, ext)
            raise HTTPException(status_code=422, detail="remplissage échoué (gabarit ou données invalides)")

        if target != ext:
            try:
                out_path = await _convert_serialized(workdir, filled_path, target)
            except asyncio.TimeoutError:
                raise HTTPException(status_code=504, detail="conversion trop longue")
            except subprocess.TimeoutExpired:
                raise HTTPException(status_code=504, detail="conversion trop longue")
            except RuntimeError:
                log.warning("tenant=%s conversion post-remplissage échouée -> %s", tenant, target)
                raise HTTPException(status_code=422, detail="conversion du document rempli échouée")
        else:
            out_path = filled_path

        data_out = open(out_path, "rb").read()
    finally:
        shutil.rmtree(workdir, ignore_errors=True)

    log.info("tenant=%s fill %s -> %s en %.1fs", tenant, ext, target, time.monotonic() - started)
    return StreamingResponse(
        io.BytesIO(data_out),
        media_type=MIME.get(target, "application/octet-stream"),
        headers={"Content-Disposition": f'attachment; filename="{stem}.{target}"'},
    )
